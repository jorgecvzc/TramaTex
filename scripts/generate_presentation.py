import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configuración de Colores (Design System)
COLOR_PRIMARY_YELLOW = RGBColor(230, 184, 0)   # #E6B800
COLOR_SECONDARY_BLUE = RGBColor(0, 35, 149)   # #002395
COLOR_BACKGROUND = RGBColor(241, 245, 249)    # #F1F5F9
COLOR_TEXT = RGBColor(30, 41, 59)             # #1E293B
COLOR_WHITE = RGBColor(255, 255, 255)

def create_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_to_slide(slide, title_text, is_primary_bg=False):
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_tf = title_shape.text_frame
    p = title_tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_PRIMARY_YELLOW if is_primary_bg else COLOR_SECONDARY_BLUE
    # Add a underline bar
    if not is_primary_bg:
        left = title_shape.left
        top = title_shape.top + title_shape.height - Inches(0.1)
        width = title_shape.width / 3
        height = Inches(0.05)
        bar = slide.shapes.add_shape(1, left, top, width, height) # 1 is rectangle
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY_YELLOW
        bar.line.fill.background()

def add_content_to_slide(slide, lines):
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for line in lines:
        line = line.strip()
        if not line: continue
        
        # Clean Marp/HTML tags
        line = re.sub(r'<.*?>', '', line)
        line = line.replace('**', '').replace('###', '').replace('##', '')
        
        if line.startswith('- '):
            p = tf.add_paragraph()
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXT
        elif line.startswith('  - '):
            p = tf.add_paragraph()
            p.text = line[4:]
            p.level = 1
            p.font.size = Pt(18)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.italic = True if 'Valor Añadido' in line or 'El Desafío' in line else False

def generate_pptx():
    prs = Presentation()
    md_path = "docs/presentations/tramatex-presentation.md"
    
    if not os.path.exists(md_path):
        print(f"Error: {md_path} no encontrado.")
        return

    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    # Split slides by Marp delimiter
    slides_raw = content.split('---')
    
    for i, slide_text in enumerate(slides_raw):
        if i == 0: continue # Skip frontmatter
        
        lines = slide_text.strip().split('\n')
        is_primary_bg = 'bg-primary' in slide_text
        
        # Determine title
        title = ""
        content_lines = []
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
            elif line.startswith('<!--') or line.startswith('header:') or line.startswith('footer:'):
                continue
            else:
                content_lines.append(line)
        
        # Create slide
        if i == 1 or i == len(slides_raw) - 1: # Portada o Final
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            create_slide_background(slide, COLOR_SECONDARY_BLUE)
            
            title_shape = slide.shapes.title
            title_shape.text = title.replace('<span class="brand">', '').replace('</span>', '')
            title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY_YELLOW
            title_shape.text_frame.paragraphs[0].font.size = Pt(60)
            
            subtitle_shape = slide.placeholders[1]
            subtitle_text = "\n".join([l for l in content_lines if l.strip() and not l.startswith('###')])
            subtitle_shape.text = subtitle_text.replace('- ', '').strip()
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            create_slide_background(slide, COLOR_BACKGROUND)
            add_title_to_slide(slide, title)
            add_content_to_slide(slide, content_lines)

    output_pptx = "docs/presentations/TramaTex_Presentacion_Final.pptx"
    prs.save(output_pptx)
    print(f"Presentación generada exitosamente en: {output_pptx}")

if __name__ == "__main__":
    generate_pptx()
